from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import numpy as np
import os

def create_presentation_from_template(template_path, topics, style="business"):
    # Check if the template file exists
    if not os.path.exists(template_path):
        print(f"Template not found at {template_path}")
        return
    
    # Load the PowerPoint template
    prs = Presentation(template_path)
    
    # Define color schemes and fonts based on style
    styles = {
        "business": {"title_color": RGBColor(0, 0, 102), "content_color": RGBColor(51, 51, 51), "background": RGBColor(240, 240, 240)},
        "project": {"title_color": RGBColor(0, 102, 204), "content_color": RGBColor(0, 51, 102), "background": RGBColor(220, 220, 220)},
        "product": {"title_color": RGBColor(255, 69, 0), "content_color": RGBColor(50, 50, 50), "background": RGBColor(255, 228, 181)},
    }
    
    chosen_style = styles.get(style, styles["business"])  # Default to business if no style is specified
    
    for index, topic in enumerate(topics):
        # Use a predefined slide layout from the template
        slide_layout = prs.slide_layouts[1]  # 1 is usually Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title for the slide with styling based on selected style
        title = slide.shapes.title
        title.text = topic['title']
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = chosen_style["title_color"]
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Align title to the left for better visual balance
        
        # Try to get the content placeholder
        try:
            content = slide.placeholders[1]  # Placeholder for content
            content.text = topic['content']
            content.text_frame.paragraphs[0].font.size = Pt(20)
            content.text_frame.paragraphs[0].font.color.rgb = chosen_style["content_color"]
            content.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Align text to the left
        except IndexError:
            # If the content placeholder doesn't exist, create a new textbox
            print("Placeholder not found, creating a new textbox.")
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(2)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = topic['content']
            text_frame.paragraphs[0].font.size = Pt(20)
            text_frame.paragraphs[0].font.color.rgb = chosen_style["content_color"]
            text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Align content left
        
        # Predefined position for the graph
        left = Inches(5.5)  # Static left position for the graph
        top = Inches(2)
        
        # Add a graph for specific slides
        if index == 1 and style in ["business", "project", "product"]:  # Adding a graph
            create_graph('graph.png', style)
            slide.shapes.add_picture('graph.png', left, top, width=Inches(4), height=Inches(3))

        # Add shapes or elements depending on the style (e.g., business uses clean lines)
        if style == "business":
            add_business_elements(slide)
        elif style == "project":
            add_project_elements(slide)
        elif style == "product":
            add_product_elements(slide)
    
    # Save the presentation
    prs.save(f'presentation_with_{style}_template.pptx')
    print(f"Presentation created successfully with {style} template!")

def create_graph(filename, style):
    # Sample data
    x = np.array(['Python', 'Java', 'C++', 'JavaScript'])
    y = np.array([4, 3, 2, 5])
    
    # Set colors based on style
    if style == "business":
        colors = ['#003366', '#0066CC', '#3399FF', '#99CCFF']
    elif style == "project":
        colors = ['#3366FF', '#3399FF', '#33CCFF', '#66FFFF']
    else:  # product
        colors = ['#FF4500', '#FFD700', '#32CD32', '#87CEEB']
    
    # Create a bar graph with colors matching the style
    plt.bar(x, y, color=colors)
    plt.title('Popularity of Programming Languages', fontsize=16)
    plt.xlabel('Languages', fontsize=12)
    plt.ylabel('Popularity Score', fontsize=12)
    
    # Tweak layout for neatness
    plt.tight_layout()
    
    # Save the graph as an image with transparency
    plt.savefig(filename, transparent=True)
    plt.close()

def add_business_elements(slide):
    # Add design elements to the corners dynamically, avoiding text overlap
    bottom_left = (Inches(0.5), Inches(6.5))  # Bottom-left corner
    top_right = (Inches(8.5), Inches(0.5))  # Top-right corner

    # Adding a bottom-left design element
    shape1 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, *bottom_left, Inches(1), Inches(0.5))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(0, 0, 102)
    
    # Adding a top-right design element
    shape2 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, *top_right, Inches(1), Inches(0.5))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0, 0, 102)

def add_project_elements(slide):
    bottom_left = (Inches(0.5), Inches(6.5))  # Bottom-left corner
    top_right = (Inches(8.5), Inches(0.5))  # Top-right corner

    # Adding a bottom-left design element
    shape1 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, *bottom_left, Inches(1), Inches(0.5))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(0, 102, 204)
    
    # Adding a top-right design element
    shape2 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, *top_right, Inches(1), Inches(0.5))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0, 102, 204)

def add_product_elements(slide):
    bottom_left = (Inches(0.5), Inches(6.5))  # Bottom-left corner
    top_right = (Inches(8.5), Inches(0.5))  # Top-right corner

    # Adding a bottom-left design element
    shape1 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, *bottom_left, Inches(1), Inches(0.5))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(255, 69, 0)
    
    # Adding a top-right design element
    shape2 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, *top_right, Inches(1), Inches(0.5))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(255, 69, 0)

if __name__ == "__main__":
    # Define static topics
    topics = [
        {'title': 'Introduction to Python', 'content': 'Python is a high-level programming language.'},
        {'title': 'Features of Python', 'content': '1. Easy to Learn\n2. Interpreted Language\n3. Extensive Libraries'},
        {'title': 'Applications of Python', 'content': '1. Web Development\n2. Data Analysis\n3. Machine Learning'},
        {'title': 'Conclusion', 'content': 'Python is versatile and powerful for various applications.'},
    ]
    
    # Full path to your PowerPoint template
    template_path = r'C:\Users\migavel\Downloads\olivedeus\template.pptx'

    # Ask user for presentation style (business, project, product)
    chosen_style = input("Choose a presentation style (business, project, product): ").lower()
    
    create_presentation_from_template(template_path, topics, style=chosen_style)
